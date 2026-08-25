#
#  Copyright 2025 The InfiniFlow Authors. All Rights Reserved.
#
#  Licensed under the Apache License, Version 2.0 (the "License");
#  you may not use this file except in compliance with the License.
#  You may obtain a copy of the License at
#
#      http://www.apache.org/licenses/LICENSE-2.0
#
#  Unless required by applicable law or agreed to in writing, software
#  distributed under the License is distributed on an "AS IS" BASIS,
#  WITHOUT WARRANTIES OR CONDITIONS OF ANY KIND, either express or implied.
#  See the License for the specific language governing permissions and
#  limitations under the License.
#
"""ppt 解析器:逐页提取 PPT 形状文本,输出每页一段拼接文本。

移植自官方 RAGFlow v0.24.0 deepdoc/parser/ppt_parser.py(逻辑 1:1)。

核心思路:
- 形状先按 (top//10, left) 排序近似还原阅读顺序(top 拿整数除法分桶,
  同一水平线的形状再按 left 排);排序结果用 id(shapes) 做实例级缓存。
- 文本框:非空段落拼起来,带项目符号的段落加 "缩进." 前缀表达层级。
- 表格(shape_type == 19):首行当表头,数据行拼成 "表头: 值;..."。
- 组合形状(shape_type == 6):递归展开内部形状。
"""

import logging
from io import BytesIO

from pptx import Presentation


class RAGFlowPptParser:
    def __init__(self):
        super().__init__()
        self._shape_cache = {}

    def __sort_shapes(self, shapes):
        cache_key = id(shapes)
        if cache_key not in self._shape_cache:
            self._shape_cache[cache_key] = sorted(
                shapes,
                key=lambda x: ((x.top if x.top is not None else 0) // 10, x.left if x.left is not None else 0)
            )
        return self._shape_cache[cache_key]

    def __get_bulleted_text(self, paragraph):
        is_bulleted = bool(paragraph._p.xpath("./a:pPr/a:buChar")) or bool(paragraph._p.xpath("./a:pPr/a:buAutoNum")) or bool(paragraph._p.xpath("./a:pPr/a:buBlip"))
        if is_bulleted:
            return f"{'  '* paragraph.level}.{paragraph.text}"
        else:
            return paragraph.text

    def __extract(self, shape):
        try:
            # First try to get text content
            if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                text_frame = shape.text_frame
                texts = []
                for paragraph in text_frame.paragraphs:
                    if paragraph.text.strip():
                        texts.append(self.__get_bulleted_text(paragraph))
                return "\n".join(texts)

            # Safely get shape_type
            try:
                shape_type = shape.shape_type
            except NotImplementedError:
                # If shape_type is not available, try to get text content
                if hasattr(shape, 'text'):
                    return shape.text.strip()
                return ""

            # Handle table
            if shape_type == 19:
                tb = shape.table
                rows = []
                for i in range(1, len(tb.rows)):
                    rows.append("; ".join([tb.cell(
                        0, j).text + ": " + tb.cell(i, j).text for j in range(len(tb.columns)) if tb.cell(i, j)]))
                return "\n".join(rows)

            # Handle group shape
            if shape_type == 6:
                texts = []
                for p in self.__sort_shapes(shape.shapes):
                    t = self.__extract(p)
                    if t:
                        texts.append(t)
                return "\n".join(texts)

            return ""

        except Exception as e:
            logging.error(f"Error processing shape: {str(e)}")
            return ""

    def __call__(self, fnm, from_page, to_page, callback=None):
        """入口:fnm 可以是文件路径(str)或二进制内容;返回每页文本组成的列表。"""
        ppt = Presentation(fnm) if isinstance(
            fnm, str) else Presentation(
            BytesIO(fnm))
        txts = []
        self.total_page = len(ppt.slides)
        for i, slide in enumerate(ppt.slides):
            if i < from_page:
                continue
            if i >= to_page:
                break
            texts = []
            for shape in self.__sort_shapes(slide.shapes):
                txt = self.__extract(shape)
                if txt:
                    texts.append(txt)
            txts.append("\n".join(texts))

        return txts
